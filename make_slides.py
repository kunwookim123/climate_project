from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# ===== 파일 경로 설정 =====
base_path = r"C:\Users\UserK\Documents\GitHub\climate_project\data\slides"
image_files = [
    "비장마_2020_02_11.png", "비장마_2021_03_20.png", "비장마_2022_09_03.png",
    "비장마_2023_11_27.png", "비장마_2024_10_14.png",
    "장마_2020_07_13.png", "장마_2021_07_03.png", "장마_2022_07_09.png",
    "장마_2023_07_18.png", "장마_2024_06_29.png"
]

output_path = os.path.join(base_path, "강수량_발전량_비교_슬라이드.pptx")

# ===== 프레젠테이션 객체 생성 =====
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ===== 제목 슬라이드 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
title_tf = title_box.text_frame
title_tf.text = "강수량과 예측 발전량 비교 (장마철 vs 비장마철)"
title_tf.paragraphs[0].font.size = Pt(40)
title_tf.paragraphs[0].font.bold = True
title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ===== 이미지별 슬라이드 생성 =====
for img_name in image_files:
    img_path = os.path.join(base_path, img_name)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 파일명에서 제목 생성
    title = img_name.replace(".png", "").replace("_", " ")

    # 제목 텍스트 상자
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(28)
    title_tf.paragraphs[0].font.bold = True

    # 이미지 추가
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.2), height=Inches(6))

# ===== 파일 저장 =====
prs.save(output_path)
print(f"✅ PPT 생성 완료: {output_path}")
